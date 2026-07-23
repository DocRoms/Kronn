"""FastAPI server — Kronn docs sidecar.

Endpoints
---------
POST /pdf    — HTML → PDF via WeasyPrint.
POST /docx   — HTML → DOCX via python-docx (heading / paragraph / list
               / table mapping; advanced CSS is lossy by design).
POST /xlsx   — structured JSON (sheets × rows) → XLSX via XlsxWriter.
POST /csv    — structured JSON (rows × cols) → CSV via stdlib.
POST /pptx   — structured JSON (slides × content) → PPTX via python-pptx.

Design notes
------------
* Loopback only (127.0.0.1). The Rust backend spawns this process with
  a discovered free port and passes it via the `KRONN_DOCS_PORT` env
  var. External callers cannot hit this port (bound address + firewall
  out-of-scope: trust the loopback boundary).
* The backend tells us WHERE to write the output file via `output_path`
  in every request body — we never choose the path ourselves. Keeps
  file lifecycle (permissions, cleanup, gc) entirely in Rust.
* On startup we print `KRONN_DOCS_READY <port>` to stdout so the Rust
  side can wait for that marker instead of polling a health endpoint —
  shaves ~80ms off cold starts. If the marker never prints within 5s
  the backend assumes the sidecar crashed (e.g. missing libs).
"""
from __future__ import annotations

import csv as _csv
import io
import logging
import os
import re
import sys
from pathlib import Path
from typing import Any, List, Optional

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, Field

logger = logging.getLogger("kronn_docs")
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")


def _configure_bundled_fontconfig() -> None:
    """Point frozen releases at a portable fontconfig search-path file.

    Homebrew/MSYS build machines have their own fontconfig configuration, but
    an installed Kronn app must not depend on those build-host paths.
    """
    if not getattr(sys, "frozen", False):
        return
    bundle_root = Path(getattr(sys, "_MEIPASS", Path(sys.executable).parent))
    config = bundle_root / "kronn_docs" / "fonts.conf"
    if config.is_file():
        os.environ["FONTCONFIG_FILE"] = str(config)


_configure_bundled_fontconfig()

app = FastAPI(title="kronn-docs", version="0.1.0")


# ─── Request / response schemas ─────────────────────────────────────────────


class PdfRequest(BaseModel):
    """Input for POST /pdf — the Rust backend decides the output path so
    every file produced by this sidecar lands where the backend expects."""
    html: str = Field(..., description="Full HTML document (can include <style>).")
    output_path: str = Field(..., description="Absolute path where the PDF will be written.")
    base_url: Optional[str] = Field(
        None,
        description="Base URL for resolving relative resources (images, CSS) referenced in the HTML. "
                    "Usually a file:// pointing at the discussion working dir.",
    )
    page_size: Optional[str] = Field(
        None,
        description="CSS @page size override (e.g. 'A4', 'Letter', '210mm 297mm'). "
                    "When unset WeasyPrint respects the HTML's own @page rules or defaults to A4.",
    )


class DocResponse(BaseModel):
    """Common response for every endpoint — path the backend can hand
    back to the UI and size_bytes for the download UX."""
    path: str
    size_bytes: int


def _pdf_page_stylesheet(html: str, page_size: Optional[str]) -> str:
    """Remove implicit page margins without trampling author margins."""
    page_rules = re.findall(r"@page[^{]*\{([^}]*)\}", html or "", flags=re.I | re.S)
    author_sets_margin = any(
        re.search(r"\bmargin(?:-(?:top|right|bottom|left))?\s*:", rule, flags=re.I)
        for rule in page_rules
    )
    declarations = []
    if page_size:
        declarations.append(f"size: {page_size}")
    if not author_sets_margin:
        declarations.append("margin: 0")
    return f"@page {{ {'; '.join(declarations)}; }}" if declarations else ""


# ─── Endpoints ──────────────────────────────────────────────────────────────


@app.get("/health")
def health() -> dict:
    """Backend liveness check. Cheap — no external deps touched."""
    return {"ok": True, "version": "0.1.0"}


@app.post("/pdf", response_model=DocResponse)
def render_pdf(req: PdfRequest) -> DocResponse:
    """HTML → PDF. WeasyPrint's `HTML(...).write_pdf(path)` does the heavy
    lifting — font resolution, page breaking, CSS paged-media features
    (@page, counter(pages)...). CSS overrides land in a stylesheet string
    when the caller forces `page_size`."""
    # Lazy-import WeasyPrint: its import chain pulls in Pango/GObject/Cairo
    # through GI bindings which is slow (~400ms on first use). Deferring
    # keeps the /health handshake snappy.
    try:
        from weasyprint import HTML, CSS  # type: ignore
    except ImportError as e:
        raise HTTPException(
            status_code=503,
            detail=f"weasyprint not installed in the sidecar venv: {e}",
        )

    stylesheets = []
    page_stylesheet = _pdf_page_stylesheet(req.html, req.page_size)
    if page_stylesheet:
        stylesheets.append(CSS(string=page_stylesheet))

    output = Path(req.output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    try:
        HTML(string=req.html, base_url=req.base_url).write_pdf(
            target=str(output),
            stylesheets=stylesheets or None,
        )
    except Exception as e:
        # WeasyPrint's error messages are usually actionable (missing font,
        # bad @page syntax, network fetch forbidden). Bubble them up so
        # Rust can surface to the user.
        logger.exception("WeasyPrint rendering failed")
        raise HTTPException(status_code=500, detail=f"PDF rendering failed: {e}") from e

    return DocResponse(path=str(output), size_bytes=output.stat().st_size)


# ─── DOCX — HTML → Word ────────────────────────────────────────────────────
#
# Word cannot represent browser layout primitives. The exporter renders the
# HTML with WeasyPrint and places the resulting pages in a DOCX, preserving
# the preview's appearance instead of producing a misleading approximation.


class DocxRequest(BaseModel):
    """Input for POST /docx — same HTML input as /pdf for workflow parity.
    Agents produce HTML once, users choose PDF or DOCX at export time."""
    html: str = Field(..., description="HTML document (headings, paragraphs, lists, tables).")
    output_path: str = Field(..., description="Absolute path where the DOCX will be written.")


@app.post("/docx", response_model=DocResponse)
def render_docx(req: DocxRequest) -> DocResponse:
    try:
        from .html_to_docx import render_html_to_docx
    except ImportError as e:
        raise HTTPException(
            status_code=503,
            detail=(
                f"DOCX renderer dependencies are unavailable: {e}. "
                "Update or reinstall Kronn, then restart the application."
            ),
        )

    output = Path(req.output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    try:
        render_html_to_docx(
            req.html,
            output,
            page_stylesheet=_pdf_page_stylesheet(req.html, "A4"),
        )
    except Exception as e:
        logger.exception("DOCX rendering failed")
        raise HTTPException(status_code=500, detail=f"DOCX rendering failed: {e}") from e

    return DocResponse(path=str(output), size_bytes=output.stat().st_size)


# ─── XLSX — structured JSON → Excel ────────────────────────────────────────


class XlsxSheet(BaseModel):
    name: str = Field(..., description="Sheet tab name (Excel max 31 chars — truncated).")
    rows: List[List[Any]] = Field(
        ..., description="2-D array of cell values. First row is rendered as bold headers by convention."
    )


class XlsxRequest(BaseModel):
    sheets: List[XlsxSheet] = Field(..., description="One or more worksheets.")
    output_path: str = Field(..., description="Absolute path where the XLSX will be written.")


@app.post("/xlsx", response_model=DocResponse)
def render_xlsx(req: XlsxRequest) -> DocResponse:
    try:
        import xlsxwriter  # type: ignore
    except ImportError as e:
        raise HTTPException(status_code=503, detail=f"XlsxWriter not installed: {e}")

    output = Path(req.output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    wb = xlsxwriter.Workbook(str(output))
    header_fmt = wb.add_format({"bold": True, "bg_color": "#eef2f5", "border": 1})
    try:
        for sheet in req.sheets:
            # Excel sheet names are capped at 31 chars and reject [\/?*[\]:].
            safe_name = re.sub(r"[\[\]\\/?*:]", "-", sheet.name)[:31] or "Sheet"
            ws = wb.add_worksheet(safe_name)
            for r_idx, row in enumerate(sheet.rows):
                for c_idx, val in enumerate(row):
                    if r_idx == 0:
                        ws.write(r_idx, c_idx, val, header_fmt)
                    else:
                        ws.write(r_idx, c_idx, val)
        wb.close()
    except Exception as e:
        logger.exception("XLSX rendering failed")
        raise HTTPException(status_code=500, detail=f"XLSX rendering failed: {e}") from e

    return DocResponse(path=str(output), size_bytes=output.stat().st_size)


# ─── CSV — structured JSON → CSV ───────────────────────────────────────────


class CsvRequest(BaseModel):
    rows: List[List[Any]] = Field(..., description="2-D array of cell values.")
    output_path: str = Field(..., description="Absolute path where the CSV will be written.")
    delimiter: Optional[str] = Field(None, description="Defaults to ','. Use ';' for FR Excel.")


@app.post("/csv", response_model=DocResponse)
def render_csv(req: CsvRequest) -> DocResponse:
    output = Path(req.output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    delim = req.delimiter if req.delimiter else ","
    try:
        with output.open("w", newline="", encoding="utf-8") as f:
            writer = _csv.writer(f, delimiter=delim)
            for row in req.rows:
                writer.writerow(row)
    except Exception as e:
        logger.exception("CSV rendering failed")
        raise HTTPException(status_code=500, detail=f"CSV rendering failed: {e}") from e

    return DocResponse(path=str(output), size_bytes=output.stat().st_size)


# ─── PPTX — structured JSON → PowerPoint ───────────────────────────────────


class PptxSlide(BaseModel):
    title: Optional[str] = Field(None, description="Slide title.")
    # Freeform content — we support either a plain `content` string (newlines
    # become bullet separators) or an explicit `bullets` list for clarity.
    content: Optional[str] = None
    bullets: Optional[List[str]] = None


class PptxRequest(BaseModel):
    slides: List[PptxSlide] = Field(..., description="One entry per slide.")
    output_path: str = Field(..., description="Absolute path where the PPTX will be written.")


@app.post("/pptx", response_model=DocResponse)
def render_pptx(req: PptxRequest) -> DocResponse:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except ImportError as e:
        raise HTTPException(
            status_code=503,
            detail=f"python-pptx not installed: {e}. Re-run `make docs-setup`.",
        )
    # Unused Inches import is fine — reserved for future layout work.
    _ = Inches

    output = Path(req.output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    # Layouts: 1 = Title + Content (header + bullets body). Fixed index
    # is stable across python-pptx versions for the default template.
    body_layout = prs.slide_layouts[1]

    try:
        for slide_spec in req.slides:
            slide = prs.slides.add_slide(body_layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = slide_spec.title or ""
            # Find the body placeholder (index 1 in the Title+Content layout).
            body = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    body = ph
                    break
            if body is None:
                continue
            tf = body.text_frame
            # Build the bullet list — prefer `bullets` (explicit) over
            # newline-split `content`.
            bullets = list(slide_spec.bullets or [])
            if not bullets and slide_spec.content:
                bullets = [b for b in slide_spec.content.split("\n") if b.strip()]
            if bullets:
                tf.text = bullets[0]
                for extra in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = extra
                    p.font.size = Pt(18)
        prs.save(str(output))
    except Exception as e:
        logger.exception("PPTX rendering failed")
        raise HTTPException(status_code=500, detail=f"PPTX rendering failed: {e}") from e

    return DocResponse(path=str(output), size_bytes=output.stat().st_size)


# ─── Entrypoint ─────────────────────────────────────────────────────────────


def main() -> None:
    """Entrypoint invoked by `python -m kronn_docs.server`.

    Reads `KRONN_DOCS_PORT` from the env (set by the Rust spawner to a
    free port it allocated). Prints `KRONN_DOCS_READY <port>` once uvicorn
    has bound and is ready to accept connections, so the Rust side can
    wait on that marker deterministically instead of polling /health.
    """
    import uvicorn
    port_str = os.environ.get("KRONN_DOCS_PORT")
    if not port_str:
        print("ERROR: KRONN_DOCS_PORT env var not set", file=sys.stderr)
        sys.exit(2)
    port = int(port_str)

    # uvicorn won't print anything useful for the ready marker, so we
    # hook the `lifespan` startup event to write it ourselves.
    @app.on_event("startup")  # type: ignore[deprecated]
    async def _announce_ready() -> None:
        # stdout is unbuffered because we spawn with PIPE'd stdio + call
        # `.reconfigure(line_buffering=True)` below. The Rust parent reads
        # line-by-line.
        print(f"KRONN_DOCS_READY {port}", flush=True)

    # Make stdout unbuffered so the ready marker lands immediately.
    try:
        sys.stdout.reconfigure(line_buffering=True)  # type: ignore[attr-defined]
    except Exception:
        pass

    uvicorn.run(
        app,
        host="127.0.0.1",
        port=port,
        log_level="info",
        access_log=False,
    )


if __name__ == "__main__":
    main()
